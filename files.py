import os
import itertools
from xlwt import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from utils import summary_plot, summary_time_plot


def get_test_file(id):
    with open(f"data/mtVRP{id}.txt") as f:
        lines = [line.rstrip() for line in f]
        lines = [line.split() for line in lines]
        lines = [[int(x) for x in line] for line in lines]
    return [lines[0], lines[1:]]


class ExcelBook:
    output_path = "./results/"

    def __init__(self, title):
        self.title = title
        self.wb = Workbook()

    def get_sheet_by_name(self, name):
        try:
            for idx in itertools.count():
                sheet = self.wb.get_sheet(idx)
                if sheet.name == name:
                    return sheet
        except IndexError:
            return self.wb.add_sheet(name)

    def add_sheet(self, index, problem_result, Th):
        sheet1 = self.get_sheet_by_name("mtVRP" + str(index))
        paths, distances, total_time = problem_result
        R = len(paths)
        for i in range(0, R):
            size = len(paths[i])
            for j in range(size):
                sheet1.write(i, j, paths[i][j])
            sheet1.write(i, size, round(distances[i], 2))
            sheet1.write(i, size + 1, 1 if distances[i] > Th else 0)

        sheet1.write(R, 0, round(sum(distances), 2))
        sheet1.write(R, 1, round(total_time, 2))
        sheet1.write(R, 2, 1)

    def save(self):
        if not os.path.exists(ExcelBook.output_path):
            os.makedirs(ExcelBook.output_path)
        self.wb.save(ExcelBook.output_path + self.title)
        print(f"{self.title} saved!")


class Slides:
    output_path = "./outputs/"

    def __init__(self, title, subtitle, total_pages=10):
        self.prs = Presentation()
        self.title_slide_layout = self.prs.slide_layouts[0]
        self.image_slide_layout = self.prs.slide_layouts[5]
        self.page = 1
        self.total_pages = total_pages
        self.data = {}

        slide = self.prs.slides.add_slide(self.title_slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(179, 229, 252)

        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True

        slide.placeholders[1].text = subtitle
        slide.placeholders[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(
            0, 0, 0
        )
        slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(16)

    def add_method_slide(self, method, file_id, description=""):
        slide = self.prs.slides.add_slide(self.image_slide_layout)
        filename = f"{method}-mtVRP{file_id}.png"

        slide.shapes.title.text = method + " results"
        slide.shapes.title.bold = True
        slide.shapes.add_picture(
            Slides.output_path + filename,
            left=Inches(1),
            top=Inches(1.5),
            width=Inches(8),
            height=Inches(5.25),
        )

        txBox = slide.shapes.add_textbox(
            left=Inches(0), top=Inches(6.4), width=Inches(5), height=Inches(1)
        )
        p = txBox.text_frame.add_paragraph()
        p.text = description
        p.font.size = Pt(14)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

        txBox = slide.shapes.add_textbox(
            left=Inches(9), top=Inches(6.75), width=Inches(1), height=Inches(1)
        )
        p = txBox.text_frame.add_paragraph()
        p.text = str("{}/{}".format(self.page, self.total_pages))
        p.font.size = Pt(10)

        self.page += 1

    def add_method_value(self, method, filename, value):
        if not method in self.data:
            self.data[method] = []
        self.data[method].append([filename, value[0], value[1]])

    def generate_timeplot(self):
        labels = []
        total_times = {}
        if "VND" in self.data:
            labels = [item[0] for item in self.data["VND"]]
            total_times["VND"] = [item[2] for item in self.data["VND"]]
        if "MS_ILS" in self.data:
            labels = [item[0] for item in self.data["MS_ILS"]]
            total_times["MS_ILS"] = [item[2] for item in self.data["MS_ILS"]]
        summary_time_plot("time_summary", labels, total_times)

        slide = self.prs.slides.add_slide(self.image_slide_layout)

        slide.shapes.title.text = "Time comparison"
        slide.shapes.title.bold = True
        slide.shapes.add_picture(
            Slides.output_path + "time_summary.png",
            left=Inches(0.25),
            top=Inches(1.5),
            width=Inches(9.5),
            height=Inches(5.5),
        )

        txBox = slide.shapes.add_textbox(
            left=Inches(9), top=Inches(6.75), width=Inches(1), height=Inches(1)
        )
        p = txBox.text_frame.add_paragraph()
        p.text = str("{}/{}".format(self.page, self.total_pages))
        p.font.size = Pt(10)

        self.page += 1

    def generate_summary(self):
        summary_plot("summary", self.data)
        slide = self.prs.slides.add_slide(self.image_slide_layout)

        slide.shapes.title.text = "Solution comparison"
        slide.shapes.title.bold = True
        slide.shapes.add_picture(
            Slides.output_path + "summary.png",
            left=Inches(0.25),
            top=Inches(1.5),
            width=Inches(9.5),
            height=Inches(5.5),
        )

        txBox = slide.shapes.add_textbox(
            left=Inches(9), top=Inches(6.75), width=Inches(1), height=Inches(1)
        )
        p = txBox.text_frame.add_paragraph()
        p.text = str("{}/{}".format(self.page, self.total_pages))
        p.font.size = Pt(10)

        self.page += 1

    def save(self, filename):
        self.prs.save(filename)
